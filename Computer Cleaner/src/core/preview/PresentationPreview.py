from __future__ import annotations

from pathlib import Path

from preview.LibreOfficePreview import build_libreoffice_thumbnail


def build_presentation_preview(path: Path, *, preview_dir: Path, thumbnail_dir: Path, max_chars: int = 4000) -> Path | None:
    thumbnail = build_libreoffice_thumbnail(path, out_dir=thumbnail_dir)
    if thumbnail is not None:
        return thumbnail

    if path.suffix.lower() != ".pptx":
        return None
    return _build_pptx_text_preview(path, out_dir=preview_dir, max_chars=max_chars)


def _build_pptx_text_preview(path: Path, *, out_dir: Path, max_chars: int) -> Path | None:
    try:
        from pptx import Presentation  # type: ignore
    except Exception:
        return None

    out_dir.mkdir(parents=True, exist_ok=True)
    out_path = out_dir / f"{path.stem}.pptx.preview.txt"
    try:
        deck = Presentation(str(path))
        chunks: list[str] = []
        for slide_index, slide in enumerate(deck.slides, start=1):
            slide_text: list[str] = []
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if text and text.strip():
                    slide_text.append(text.strip())
            if slide_text:
                chunks.append(f"Slide {slide_index}\n" + "\n".join(slide_text))

        content = "\n\n".join(chunks).strip()
        if not content:
            return None
        out_path.write_text(content[:max_chars], encoding="utf-8", errors="replace")
        return out_path
    except Exception:
        return None
